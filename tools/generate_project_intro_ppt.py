#!/usr/bin/env python3
"""Generate a PPTX deck from the project intro markdown outline."""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_INPUT = "work_scr/new_gui_project_intro_ppt.md"
DEFAULT_OUTPUT = "work_scr/new_gui_project_intro.pptx"


def strip_markdown(text: str) -> str:
    """Remove simple markdown markers for PPT rendering."""
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    return text.strip()


def parse_slides(markdown_text: str) -> list[dict[str, object]]:
    """Parse slide data from the markdown source."""
    sections = [section.strip() for section in markdown_text.split("\n---\n") if section.strip()]
    slides = []

    for section in sections:
        lines = [line.rstrip() for line in section.splitlines()]
        title = ""
        bullets: list[str] = []
        paragraphs: list[str] = []

        for line in lines:
            stripped = line.strip()
            if not stripped or stripped.startswith("# New-GUI Project Intro PPT"):
                continue
            if stripped.startswith("## Slide"):
                continue
            if stripped.startswith("Speaker notes:"):
                break
            if stripped.startswith("- "):
                bullets.append(strip_markdown(stripped[2:]))
                continue
            cleaned = strip_markdown(stripped)
            if cleaned and not title:
                title = cleaned
            elif cleaned:
                paragraphs.append(cleaned)

        if title:
            slides.append({"title": title, "paragraphs": paragraphs, "bullets": bullets})

    return slides


def add_title_slide(prs: Presentation, slide_data: dict[str, object]) -> None:
    """Add the opening slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = str(slide_data["title"])
    subtitle = slide.placeholders[1]
    subtitle_lines = [str(line) for line in slide_data["paragraphs"]] + [f"- {item}" for item in slide_data["bullets"]]
    subtitle.text = "\n".join(subtitle_lines[:4])


def add_content_slide(prs: Presentation, slide_data: dict[str, object]) -> None:
    """Add a title-and-content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = str(slide_data["title"])

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    paragraphs = [str(item) for item in slide_data["paragraphs"]]
    bullets = [str(item) for item in slide_data["bullets"]]
    content_items = paragraphs + bullets

    if not content_items:
        content_items = [" "]

    for index, item in enumerate(content_items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(20 if index == 0 and not bullets else 18)
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT


def add_footer(prs: Presentation) -> None:
    """Add a small footer to all slides except the title slide."""
    for index, slide in enumerate(prs.slides):
        if index == 0:
            continue
        textbox = slide.shapes.add_textbox(Inches(0.3), Inches(6.9), Inches(12.5), Inches(0.3))
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.text = "New-GUI Project Intro"
        paragraph.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.RIGHT


def build_presentation(slides: list[dict[str, object]]) -> Presentation:
    """Build the presentation object."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if not slides:
        raise SystemExit("No slides parsed from markdown source.")

    add_title_slide(prs, slides[0])
    for slide_data in slides[1:]:
        add_content_slide(prs, slide_data)
    add_footer(prs)
    return prs


def parse_args() -> argparse.Namespace:
    """Parse CLI arguments."""
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", default=DEFAULT_INPUT, help="Markdown source file")
    parser.add_argument("--output", default=DEFAULT_OUTPUT, help="Output PPTX file")
    return parser.parse_args()


def main() -> int:
    """CLI entrypoint."""
    args = parse_args()
    repo_root = Path(__file__).resolve().parent.parent
    input_path = (repo_root / args.input).resolve()
    output_path = (repo_root / args.output).resolve()

    markdown_text = input_path.read_text(encoding="utf-8")
    slides = parse_slides(markdown_text)
    presentation = build_presentation(slides)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    print(output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
